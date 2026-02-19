"""
Generates project14-poster.pptx — A0 portrait poster for Project 14.

Run:  python project/make_poster.py
Output: project/project14-poster.pptx

Edit the CONTENT dict below to fill in paper details as you read.
"""

from pptx import Presentation
from pptx.util import Mm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
import copy

# ──────────────────────────────────────────────────────────────────────────────
# CONTENT — edit this dict as you read the paper
# ──────────────────────────────────────────────────────────────────────────────
CONTENT = {
    "title": (
        "A Multi-Level Interpretable Sleep Stage Scoring System\n"
        "by Infusing Experts' Knowledge Into a Deep Network Architecture"
    ),
    "authors": "H. Niknazar  &  S. C. Mednick",
    "journal": "IEEE Trans. Pattern Analysis and Machine Intelligence, vol. 46, no. 7, pp. 5044–5061, July 2024",
    "presented_by": "Presented by: [Your Name] — Multimodal Systems, UniGe 2025–2026",

    # ── MOTIVATION & OBJECTIVES ───────────────────────────────────────────────
    "motivation": [
        "Sleep staging: experts score 30-second PSG epochs\n"
        "into Wake / N1 / N2 / N3 / REM following AASM rules.\n"
        "Gold standard, but expensive and slow.",

        "Inter-rater agreement ~80% — even experts disagree,\n"
        "especially on N1 (the most transitional stage).",

        "Deep learning systems match human accuracy\n"
        "but are black boxes — clinicians cannot inspect,\n"
        "verify, or correct their decisions.",

        "GOAL: A system that is accurate AND interpretable\n"
        "by systematically injecting AASM expert knowledge\n"
        "into the network at multiple levels of abstraction.",
    ],
    "course_connection": (
        "MM-Systems connection:\n"
        "EEG + EOG + EMG = 3 physiological modalities\n"
        "→ multimodal system (W3C / Jaimes-Sebe definition)\n\n"
        "Architecture maps to the four-layer hierarchy\n"
        "(Lecture 3): physical signals → low-level features\n"
        "→ mid-level features → concepts (sleep stage label)\n\n"
        "Complementarity (Lecture 2): EEG carries brain-state\n"
        "info; EOG is essential for REM; EMG for atonia."
    ),

    # ── METHODOLOGY ───────────────────────────────────────────────────────────
    "input_modalities": (
        "Input — polysomnography (PSG), 30-second epochs:\n"
        "• EEG: frontal (F3/F4), central (C3/C4), occipital (O1/O2)\n"
        "  derivations per AASM — spindles, K-complexes, slow waves\n"
        "• EOG: left + right eye (E1-M2, E2-M1)\n"
        "  — rapid eye movements detect REM\n"
        "• EMG: submental / chin muscle\n"
        "  — atonia distinguishes REM from N1"
    ),
    "architecture_levels": [
        "Level 1 — Raw PSG signal\n(EEG + EOG + EMG epochs @ 100–256 Hz)",
        "Level 2 — Time/frequency features\n(CNN extracts spindles, K-complexes,\nslow waves, band power)",
        "Level 3 — Expert-rule representations\n(attention guided by AASM criteria;\nauxiliary losses align with expert features)",
        "Level 4 — Sleep stage label\n(Wake / N1 / N2 / N3 / REM\n+ sequence transition constraints)",
    ],
    "expert_injection": (
        "Multi-mechanism knowledge infusion:\n\n"
        "1. Attention regularisation — weights guided toward\n"
        "   AASM-relevant time/frequency regions per stage:\n"
        "   σ-band (12–15 Hz) → N2 spindles\n"
        "   δ-band (0.5–4 Hz) → N3 slow waves\n"
        "   rapid EOG bursts → REM\n\n"
        "2. Auxiliary losses — soft expert-rule constraints:\n"
        "   L = L_CE + λ₁·L_feat + λ₂·L_trans + λ₃·L_attn\n\n"
        "3. Transition priors — expert-derived sleep-stage\n"
        "   transition matrix constrains the sequence model\n"
        "   (e.g., N3 does not directly follow REM)"
    ),
    "datasets": (
        "Evaluation datasets:\n"
        "• SleepEDF-20 / SleepEDF-78 — healthy adults,\n"
        "  2-ch EEG (Fpz-Cz, Pz-Oz) + EOG + EMG\n"
        "• SHHS (Sleep Heart Health Study) — ~5,800 subjects,\n"
        "  1-ch EEG (C4-A1) + EOG + EMG; large-scale benchmark\n"
        "• MASS — ~200 subjects, full PSG montage\n"
        "• ISRUC — clinical population including patients\n\n"
        "Metrics: Accuracy, Macro-F1, Cohen's κ"
    ),

    # ── KEY RESULTS ───────────────────────────────────────────────────────────
    "results_table": (
        "Comparison on SleepEDF-20 (verify against paper):\n\n"
        "Method                  ACC    MF1    κ\n"
        "────────────────────────────────────────\n"
        "DeepSleepNet (2017)     82.0   76.9   0.76\n"
        "SeqSleepNet  (2019)     87.1   80.0   0.83\n"
        "SleepTransformer (2022) 88.0   81.0   0.84\n"
        "XSleepNet    (2021)     87.7   80.5   0.83\n"
        "Niknazar & Mednick      ~89    ~83    ~0.85\n\n"
        "* Figures from training-data reconstruction.\n"
        "  Replace with exact values from paper Table."
    ),
    "per_stage_f1": (
        "Per-stage F1 (SleepEDF — approximate):\n"
        "Wake  ~90% | N1  ~45–55% | N2  ~87%\n"
        "N3  ~85% | REM  ~85%\n\n"
        "N1 is hardest: transitional EEG, highest\n"
        "inter-rater disagreement (~52–60% among experts).\n"
        "Expert knowledge infusion gives largest N1 gain."
    ),
    "interpretability_result": (
        "Interpretability outputs:\n"
        "• Temporal attention maps — highlight spindles\n"
        "  (N2), slow waves (N3), sawtooth waves (REM)\n"
        "  overlaid on raw EEG: clinician-verifiable\n"
        "• Frequency-band attention — shows which\n"
        "  Hz bands drove each stage decision\n"
        "• Expert-rule fidelity metric: quantifies\n"
        "  how often attention aligns with AASM criteria\n"
        "• Transition probability matrices: comparable\n"
        "  to expert-established hypnogram statistics"
    ),

    # ── DISCUSSION ────────────────────────────────────────────────────────────
    "strengths": [
        "Clinician-verifiable: attention maps align\nwith known AASM stage markers",
        "State-of-the-art accuracy on multiple\npublic benchmarks simultaneously",
        "Multi-mechanism infusion: works at signal,\nepoch, and sequence levels",
        "N1 F1 improvement over black-box DL\n(most clinically meaningful gain)",
    ],
    "limitations": [
        "Relies on automatic spindle/K-complex detectors\nfor auxiliary supervision — noisy labels propagate",
        "AASM-rule encoding is adult-specific;\npediatric or disorder-specific rules not addressed",
        "N1 remains hard: human ceiling ~52–60%\nagreement sets fundamental limit",
        "Hyperparameters λ₁, λ₂, λ₃ require dataset-\nspecific tuning; adds training complexity",
    ],
    "future_work": (
        "• Extend knowledge injection to disorder-specific\n"
        "  rules (sleep apnea, narcolepsy, insomnia)\n"
        "• Incorporate clinician feedback in the loop\n"
        "  (active learning / human-in-the-loop)\n"
        "• Explore large-scale pre-training with\n"
        "  fine-tuning on clinical populations\n"
        "• [Verify against paper's stated future work]"
    ),
}

# ──────────────────────────────────────────────────────────────────────────────
# DESIGN CONSTANTS
# ──────────────────────────────────────────────────────────────────────────────

# A0 portrait in EMUs (1 mm = 36000 EMU)
A0_W = Mm(841)
A0_H = Mm(1189)

# Colours
C_DARK_BLUE  = RGBColor(0x1A, 0x37, 0x5E)  # deep navy
C_MID_BLUE   = RGBColor(0x2E, 0x6D, 0xA4)  # medium blue
C_LIGHT_BLUE = RGBColor(0xD6, 0xE8, 0xF5)  # pale blue fill
C_ORANGE     = RGBColor(0xE8, 0x74, 0x22)  # accent orange
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK      = RGBColor(0x11, 0x11, 0x11)
C_GREY_BG    = RGBColor(0xF4, 0xF4, 0xF4)
C_GREY_LINE  = RGBColor(0xCC, 0xCC, 0xCC)

# Margins & layout geometry (mm)
MAR   = 15   # outer margin
GAP   = 8    # gap between columns
TITLE_H = 90  # title bar height (mm)
DISC_H  = 130  # discussion bar height (mm)
COL_COUNT = 3


def mm(v): return Mm(v)


def add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_w_pt=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        mm(x), mm(y), mm(w), mm(h)
    )
    fill = shape.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    line = shape.line
    if line_rgb:
        line.color.rgb = line_rgb
        line.width = Pt(line_w_pt)
    else:
        line.fill.background()
    return shape


def add_text_box(slide, x, y, w, h, text, font_size, bold=False,
                 color=C_BLACK, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(mm(x), mm(y), mm(w), mm(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb


def add_section_header(slide, x, y, w, h, title, fill=C_MID_BLUE):
    box = add_rect(slide, x, y, w, h, fill_rgb=fill)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = C_WHITE
    run.font.name = "Calibri"
    return box


def add_content_box(slide, x, y, w, h, text, font_size=18, fill=C_GREY_BG):
    box = add_rect(slide, x, y, w, h, fill_rgb=fill, line_rgb=C_GREY_LINE, line_w_pt=0.5)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = C_BLACK
    run.font.name = "Calibri"
    return box


def add_bullet_box(slide, x, y, w, h, bullets, font_size=17, fill=C_GREY_BG):
    box = add_rect(slide, x, y, w, h, fill_rgb=fill, line_rgb=C_GREY_LINE, line_w_pt=0.5)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = "• " + bullet
        run.font.size = Pt(font_size)
        run.font.color.rgb = C_BLACK
        run.font.name = "Calibri"
    return box


def add_arch_diagram(slide, x, y, w, h, levels):
    """Draw a simple vertical flow diagram for the multi-level architecture."""
    n = len(levels)
    box_h = (h - (n - 1) * 4) / n  # height of each level box
    arrow_h = 4  # mm

    level_colors = [
        RGBColor(0x9B, 0xC2, 0xE6),  # level 1 — lightest blue
        RGBColor(0x5B, 0x9C, 0xD6),  # level 2
        RGBColor(0x2E, 0x6D, 0xA4),  # level 3
        RGBColor(0x1A, 0x37, 0x5E),  # level 4 — darkest
    ]

    for i, level_text in enumerate(levels):
        cur_y = y + i * (box_h + arrow_h)
        color = level_colors[i] if i < len(level_colors) else C_MID_BLUE
        box = add_rect(slide, x, cur_y, w, box_h, fill_rgb=color,
                       line_rgb=C_WHITE, line_w_pt=1)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = level_text
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = C_WHITE
        run.font.name = "Calibri"

        # Arrow between boxes
        if i < n - 1:
            arr_y = cur_y + box_h
            arr = add_rect(slide, x + w * 0.4, arr_y, w * 0.2, arrow_h,
                           fill_rgb=C_ORANGE)


# ──────────────────────────────────────────────────────────────────────────────
# BUILD PRESENTATION
# ──────────────────────────────────────────────────────────────────────────────

def build_poster(out_path="project/project14-poster.pptx"):
    prs = Presentation()
    prs.slide_width  = A0_W
    prs.slide_height = A0_H

    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)

    # Background
    add_rect(slide, 0, 0, 841, 1189, fill_rgb=C_WHITE)

    # ── TITLE BAR ─────────────────────────────────────────────────────────────
    add_rect(slide, 0, 0, 841, TITLE_H, fill_rgb=C_DARK_BLUE)

    add_text_box(
        slide, MAR, 8, 841 - 2*MAR, 40,
        CONTENT["title"],
        font_size=34, bold=True, color=C_WHITE,
        align=PP_ALIGN.CENTER
    )
    add_text_box(
        slide, MAR, 50, 841 - 2*MAR, 14,
        CONTENT["authors"],
        font_size=22, bold=False, color=C_LIGHT_BLUE,
        align=PP_ALIGN.CENTER
    )
    add_text_box(
        slide, MAR, 65, 841 - 2*MAR, 12,
        CONTENT["journal"],
        font_size=16, bold=False, color=C_LIGHT_BLUE,
        align=PP_ALIGN.CENTER
    )
    add_text_box(
        slide, MAR, 78, 841 - 2*MAR, 10,
        CONTENT["presented_by"],
        font_size=14, bold=False, color=RGBColor(0xAA, 0xCC, 0xEE),
        align=PP_ALIGN.CENTER
    )

    # ── THREE-COLUMN BODY ─────────────────────────────────────────────────────
    body_top  = TITLE_H + 8
    body_bot  = 1189 - DISC_H - 8
    body_h    = body_bot - body_top
    total_w   = 841 - 2*MAR - 2*GAP
    col_w     = total_w / COL_COUNT

    col_x = [MAR + i*(col_w + GAP) for i in range(COL_COUNT)]

    HDR_H = 22   # section header height (mm)
    PAD   = 4    # padding below header before content

    # ── COLUMN 1: MOTIVATION & OBJECTIVES ────────────────────────────────────
    c1x, c1y, c1w = col_x[0], body_top, col_w

    add_section_header(slide, c1x, c1y, c1w, HDR_H, "MOTIVATION & OBJECTIVES")

    # Motivation bullets
    cur_y = c1y + HDR_H + PAD
    mot_h = 90
    add_bullet_box(slide, c1x, cur_y, c1w, mot_h, CONTENT["motivation"], font_size=17)

    # Course connection
    cur_y += mot_h + PAD
    cc_h = 65
    add_section_header(slide, c1x, cur_y, c1w, HDR_H, "MM-SYSTEMS CONNECTION",
                       fill=C_ORANGE)
    cur_y += HDR_H + PAD
    add_content_box(slide, c1x, cur_y, c1w, cc_h - HDR_H - PAD,
                    CONTENT["course_connection"], font_size=16,
                    fill=RGBColor(0xFF, 0xF0, 0xE0))

    # Input modalities
    cur_y += (cc_h - HDR_H - PAD) + PAD
    add_section_header(slide, c1x, cur_y, c1w, HDR_H, "INPUT MODALITIES",
                       fill=C_MID_BLUE)
    cur_y += HDR_H + PAD
    inp_h = 60
    add_content_box(slide, c1x, cur_y, c1w, inp_h,
                    CONTENT["input_modalities"], font_size=17)

    # ── COLUMN 2: METHODOLOGY ─────────────────────────────────────────────────
    c2x, c2y, c2w = col_x[1], body_top, col_w

    add_section_header(slide, c2x, c2y, c2w, HDR_H, "METHODOLOGY")

    # Architecture diagram
    cur_y = c2y + HDR_H + PAD
    arch_h = 110
    add_arch_diagram(slide, c2x, cur_y, c2w, arch_h,
                     CONTENT["architecture_levels"])

    # Expert injection
    cur_y += arch_h + PAD
    add_section_header(slide, c2x, cur_y, c2w, HDR_H, "EXPERT KNOWLEDGE INJECTION",
                       fill=C_MID_BLUE)
    cur_y += HDR_H + PAD
    ej_h = 80
    add_content_box(slide, c2x, cur_y, c2w, ej_h,
                    CONTENT["expert_injection"], font_size=17)

    # Datasets
    cur_y += ej_h + PAD
    add_section_header(slide, c2x, cur_y, c2w, HDR_H, "DATASETS & EVALUATION",
                       fill=C_MID_BLUE)
    cur_y += HDR_H + PAD
    ds_h = 65
    add_content_box(slide, c2x, cur_y, c2w, ds_h,
                    CONTENT["datasets"], font_size=17)

    # ── COLUMN 3: KEY RESULTS ─────────────────────────────────────────────────
    c3x, c3y, c3w = col_x[2], body_top, col_w

    add_section_header(slide, c3x, c3y, c3w, HDR_H, "KEY RESULTS")

    # Results comparison table
    cur_y = c3y + HDR_H + PAD
    rt_h = 90
    add_content_box(slide, c3x, cur_y, c3w, rt_h,
                    CONTENT["results_table"], font_size=15,
                    fill=RGBColor(0xEB, 0xF3, 0xFB))

    # Per-stage F1
    cur_y += rt_h + PAD
    add_section_header(slide, c3x, cur_y, c3w, HDR_H, "PER-STAGE PERFORMANCE",
                       fill=C_MID_BLUE)
    cur_y += HDR_H + PAD
    ps_h = 45
    add_content_box(slide, c3x, cur_y, c3w, ps_h,
                    CONTENT["per_stage_f1"], font_size=17)

    # Interpretability
    cur_y += ps_h + PAD
    add_section_header(slide, c3x, cur_y, c3w, HDR_H, "INTERPRETABILITY",
                       fill=C_MID_BLUE)
    cur_y += HDR_H + PAD
    ip_h = 70
    add_content_box(slide, c3x, cur_y, c3w, ip_h,
                    CONTENT["interpretability_result"], font_size=17)

    # ── DISCUSSION BAR ────────────────────────────────────────────────────────
    disc_top = 1189 - DISC_H
    add_rect(slide, 0, disc_top - 4, 841, 4, fill_rgb=C_ORANGE)  # orange rule

    add_section_header(slide, 0, disc_top, 841, HDR_H, "DISCUSSION",
                       fill=C_DARK_BLUE)

    disc_body_top = disc_top + HDR_H + PAD
    disc_body_h   = DISC_H - HDR_H - PAD - 4
    disc_third    = (841 - 2*MAR - 2*GAP) / 3

    # Strengths
    dx = MAR
    add_section_header(slide, dx, disc_body_top, disc_third, HDR_H,
                       "Strengths", fill=RGBColor(0x27, 0x7A, 0x3E))
    add_bullet_box(slide, dx, disc_body_top + HDR_H + PAD,
                   disc_third, disc_body_h - HDR_H - PAD,
                   CONTENT["strengths"], font_size=17,
                   fill=RGBColor(0xE8, 0xF5, 0xEB))

    # Limitations
    dx = MAR + disc_third + GAP
    add_section_header(slide, dx, disc_body_top, disc_third, HDR_H,
                       "Limitations", fill=RGBColor(0xC0, 0x39, 0x2B))
    add_bullet_box(slide, dx, disc_body_top + HDR_H + PAD,
                   disc_third, disc_body_h - HDR_H - PAD,
                   CONTENT["limitations"], font_size=17,
                   fill=RGBColor(0xFD, 0xED, 0xEC))

    # Future Work
    dx = MAR + 2*(disc_third + GAP)
    add_section_header(slide, dx, disc_body_top, disc_third, HDR_H,
                       "Future Work", fill=C_MID_BLUE)
    add_content_box(slide, dx, disc_body_top + HDR_H + PAD,
                    disc_third, disc_body_h - HDR_H - PAD,
                    CONTENT["future_work"], font_size=17)

    # ── SAVE ──────────────────────────────────────────────────────────────────
    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    import os
    os.makedirs("project", exist_ok=True)
    build_poster("project/project14-poster.pptx")
